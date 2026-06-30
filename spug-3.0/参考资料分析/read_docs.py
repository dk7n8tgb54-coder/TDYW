
import os
from docx import Document
from pptx import Presentation

def read_docx(file_path):
    """读取docx文件内容"""
    doc = Document(file_path)
    content = []
    for para in doc.paragraphs:
        if para.text.strip():
            content.append(para.text)
    return '\n'.join(content)

def read_pptx(file_path):
    """读取pptx文件内容"""
    prs = Presentation(file_path)
    content = []
    for i, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text)
        if slide_text:
            content.append(f"=== 第{i+1}页 ===")
            content.extend(slide_text)
    return '\n'.join(content)

def main():
    base_dir = r"E:\TDYW\spug-3.0\参考资料分析"
    output_dir = r"E:\TDYW\spug-3.0\参考资料分析\提取内容"
    os.makedirs(output_dir, exist_ok=True)
    
    files = [
        ("28_AIOps项目申报书_正式版.docx", "docx"),
        ("33_AIOps项目申报书_终版.docx", "docx"),
        ("34_湖北烟草科技项目申报书.docx", "docx"),
        ("ITSM运维服务管理平台介绍.pptx", "pptx"),
        ("ITSM整体规划.pptx", "pptx"),
    ]
    
    for filename, filetype in files:
        file_path = os.path.join(base_dir, filename)
        if not os.path.exists(file_path):
            print(f"文件不存在: {filename}")
            continue
            
        print(f"正在读取: {filename}")
        try:
            if filetype == "docx":
                content = read_docx(file_path)
            elif filetype == "pptx":
                content = read_pptx(file_path)
            
            output_file = os.path.join(output_dir, f"{os.path.splitext(filename)[0]}.txt")
            with open(output_file, 'w', encoding='utf-8') as f:
                f.write(content)
            print(f"  已保存到: {output_file}")
            print(f"  字数: {len(content)}")
        except Exception as e:
            print(f"  读取失败: {e}")

if __name__ == "__main__":
    main()
